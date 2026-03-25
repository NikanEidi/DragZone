import io
from fastapi import UploadFile

# Parsers
from pypdf import PdfReader
from bs4 import BeautifulSoup

try:
    import docx
except ImportError:
    docx = None

try:
    import pandas as pd
except ImportError:
    pd = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

from core.llm import client

async def parse_file(file: UploadFile) -> str:
    filename = file.filename.lower()
    content_type = file.content_type
    
    try:
        file_bytes = await file.read()
        extracted_content = ""
        
        # 1. PDF Parsing
        if filename.endswith(".pdf"):
            reader = PdfReader(io.BytesIO(file_bytes))
            for page in reader.pages:
                extracted_content += page.extract_text() + "\n"
                
        # 2. Word Document Parsing (.docx)
        elif filename.endswith(".docx") or filename.endswith(".doc"):
            if docx:
                doc = docx.Document(io.BytesIO(file_bytes))
                extracted_content = "\n".join([para.text for para in doc.paragraphs])
            else:
                extracted_content = "[WARNING: python-docx not installed. Could not parse DOCX.]"

        # 3. Excel / CSV Parsing
        elif filename.endswith((".xlsx", ".xls", ".csv")):
            if pd:
                try:
                    if filename.endswith(".csv"):
                        df = pd.read_csv(io.BytesIO(file_bytes))
                    else:
                        df = pd.read_excel(io.BytesIO(file_bytes))
                    extracted_content = df.to_markdown(index=False)
                except Exception as e:
                    extracted_content = f"[ERROR PARSING EXCEL/CSV: {str(e)}]"
            else:
                extracted_content = "[WARNING: pandas not installed. Could not parse Spreadsheet.]"

        # 4. PowerPoint Parsing (.pptx)
        elif filename.endswith(".pptx"):
            if Presentation:
                try:
                    prs = Presentation(io.BytesIO(file_bytes))
                    for slide_num, slide in enumerate(prs.slides):
                        extracted_content += f"\n--- Slide {slide_num + 1} ---\n"
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                extracted_content += shape.text + "\n"
                except Exception as e:
                    extracted_content = f"[ERROR PARSING PPTX: {str(e)}]"
            else:
                extracted_content = "[WARNING: python-pptx not installed. Could not parse PPTX.]"

        # 5. HTML / XML Parsing
        elif filename.endswith((".html", ".htm", ".xml")):
            soup = BeautifulSoup(file_bytes, "html.parser" if filename.endswith(".html") else "xml")
            extracted_content = soup.get_text(separator="\n", strip=True)

        # 6. Image Analysis via LLaVA
        elif content_type and content_type.startswith("image/"):
            try:
                response = await client.generate(
                    model='llava',
                    prompt="Describe this image technically. Extract any visible text or data. If it is a diagram, explain the flow.",
                    images=[file_bytes],
                    stream=False
                )
                extracted_content = response.get('response', 'No analysis generated.')
            except Exception as e:
                extracted_content = f"[ERROR PARSING IMAGE: {str(e)}]"

        # 7. Fallback to raw text/code decoding
        else:
            try:
                extracted_content = file_bytes.decode("utf-8")
            except UnicodeDecodeError:
                extracted_content = f"[ERROR PARSING {file.filename}]: Could not decode file as UTF-8. Binary format not fully supported without specific extension."

        return f"\n[--- START FILE: {file.filename} ---]\n{extracted_content}\n[--- END FILE ---]\n"

    except Exception as e:
        return f"\n[ERROR PARSING {file.filename}]: {str(e)}\n"
